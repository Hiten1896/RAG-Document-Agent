"""Document text extraction.

Every parser returns a list of (label, text) units — a label is a page/slide/sheet
marker like "page 3" or "slide 7". Keeping units separate rather than returning one
flat string is what lets `main.py` attach a real `page` value to each chunk's
metadata, which the frontend's source cards render. A parser for a format with no
natural pagination (csv, txt) just returns a single unit.

Dispatch goes through EXTENSION_PARSERS, so adding a format is one dict entry and
SUPPORTED_EXTENSIONS stays in sync automatically.
"""

import csv
import os
from typing import Any, Callable, Dict, List, Optional, Tuple

import fitz  # PyMuPDF

# A text unit and the human-readable location it came from.
Unit = Tuple[str, str]

# Legacy binary Office formats need a LibreOffice/antiword conversion step, which
# won't fit Render's free tier. Named explicitly so users get told what to do
# instead of a bare "unsupported file type".
LEGACY_FORMAT_ADVICE = {
    '.doc': 'Save it as .docx and upload again.',
    '.ppt': 'Save it as .pptx and upload again.',
    '.xls': 'Save it as .xlsx and upload again.',
}


def load_pdf_fast(file_path: str) -> List[Unit]:
    """PDF extraction. Extracts digital text directly; no OCR fallback.

    Pages with no extractable digital text (i.e. scanned image pages) are
    simply skipped — running OCR (Tesseract) on Render's free tier was heavy
    enough to push the service past its memory/CPU budget and take the
    server offline. This project intentionally does not support scanned
    (image-only) documents.
    """
    units: List[Unit] = []
    doc = fitz.open(file_path)

    try:
        for index, page in enumerate(doc, start=1):
            page_text = page.get_text("text")

            if page_text and len(page_text.strip()) > 10:
                units.append((f"page {index}", page_text))
    finally:
        doc.close()

    return units


def load_docx(file_path: str) -> List[Unit]:
    """Extract paragraphs and table cell text from a .docx.

    python-docx exposes tables separately from paragraphs, so a document whose
    content lives mostly in tables would come back nearly empty if only
    `doc.paragraphs` were read.
    """
    from docx import Document

    doc = Document(file_path)
    parts: List[str] = [p.text for p in doc.paragraphs if p.text and p.text.strip()]

    for table in doc.tables:
        for row in table.rows:
            cells = [cell.text.strip() for cell in row.cells if cell.text and cell.text.strip()]
            if cells:
                parts.append(" | ".join(cells))

    # .docx has no addressable page count without rendering it, so the whole
    # document is one unit rather than reporting a page number we can't know.
    text = "\n".join(parts)
    return [("document", text)] if text.strip() else []


def load_pptx(file_path: str) -> List[Unit]:
    """Extract shape text and speaker notes from a .pptx, one unit per slide."""
    from pptx import Presentation

    prs = Presentation(file_path)
    units: List[Unit] = []

    for index, slide in enumerate(prs.slides, start=1):
        parts: List[str] = []

        for shape in slide.shapes:
            # Not every shape holds text (images, lines, plain rectangles).
            if not getattr(shape, "has_text_frame", False):
                continue
            shape_text = shape.text_frame.text
            if shape_text and shape_text.strip():
                parts.append(shape_text.strip())

        # Speaker notes often carry the real explanation behind a sparse slide,
        # so they're worth indexing alongside the visible text.
        if slide.has_notes_slide:
            notes = slide.notes_slide.notes_text_frame
            if notes is not None and notes.text and notes.text.strip():
                parts.append(f"Speaker notes: {notes.text.strip()}")

        if parts:
            units.append((f"slide {index}", "\n".join(parts)))

    return units


def load_xlsx(file_path: str) -> List[Unit]:
    """Extract cell values from a .xlsx, one unit per worksheet.

    read_only=True streams rows instead of building the full object model, which
    keeps memory flat on large workbooks. data_only=True yields cached formula
    results rather than the formula strings, which is what a reader wants.
    """
    from openpyxl import load_workbook

    wb = load_workbook(file_path, read_only=True, data_only=True)
    units: List[Unit] = []

    try:
        for sheet in wb.worksheets:
            rows: List[str] = []
            for row in sheet.iter_rows(values_only=True):
                cells = [str(v).strip() for v in row if v is not None and str(v).strip()]
                if cells:
                    rows.append(" | ".join(cells))
            if rows:
                units.append((f"sheet {sheet.title}", "\n".join(rows)))
    finally:
        wb.close()

    return units


def load_csv(file_path: str) -> List[Unit]:
    """Read a CSV into pipe-delimited lines."""
    rows: List[str] = []
    with open(file_path, "r", encoding="utf-8", errors="ignore", newline="") as f:
        for row in csv.reader(f):
            cells = [c.strip() for c in row if c and c.strip()]
            if cells:
                rows.append(" | ".join(cells))

    text = "\n".join(rows)
    return [("document", text)] if text.strip() else []


def load_text(file_path: str) -> List[Unit]:
    """Read plain text and markdown files."""
    with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
        text = f.read()
    return [("document", text)] if text.strip() else []


# Extension -> parser. SUPPORTED_EXTENSIONS derives from this, so a new format
# only needs a single entry here.
EXTENSION_PARSERS: Dict[str, Callable[[str], List[Unit]]] = {
    '.pdf': load_pdf_fast,
    '.docx': load_docx,
    '.pptx': load_pptx,
    '.xlsx': load_xlsx,
    '.csv': load_csv,
    '.txt': load_text,
    '.md': load_text,
}

SUPPORTED_EXTENSIONS = frozenset(EXTENSION_PARSERS)


def load_document(file_path: str, original_filename: Optional[str] = None) -> Dict[str, Any]:
    """Main document loader dispatcher.

    Returns the joined text plus the per-unit breakdown, so callers can chunk each
    unit independently and keep its page/slide label attached.
    """
    if not os.path.exists(file_path):
        raise FileNotFoundError(f"File not found: {file_path}")

    # The upload path is a temp file that may carry no extension, so prefer the
    # caller-supplied original filename when deciding how to parse.
    filename_for_ext = original_filename if original_filename else file_path
    ext = os.path.splitext(filename_for_ext)[1].lower()

    if ext in LEGACY_FORMAT_ADVICE:
        raise ValueError(
            f"{ext} is an older binary Office format this server can't read. "
            f"{LEGACY_FORMAT_ADVICE[ext]}"
        )

    parser = EXTENSION_PARSERS.get(ext)
    if parser is None:
        supported = ", ".join(sorted(SUPPORTED_EXTENSIONS))
        raise ValueError(f"Unsupported file type: {ext or '(none)'}. Supported: {supported}")

    units = parser(file_path)

    return {
        "file_path": file_path,
        "file_name": original_filename or os.path.basename(file_path),
        "extension": ext,
        "units": units,
        "unit_count": len(units),
        "content": "\n\n".join(text for _, text in units),
    }